import re
import uuid
from pathlib import Path

from bson import ObjectId
from flask import Blueprint, jsonify, request, send_file
from flask_jwt_extended import jwt_required
from pptx import Presentation

from extensions import LIBRARY_DIR, output_path, presentations_col
from models import User
from shared.auth_utils import get_current_user_doc
from shared.pptx_utils import (
    apply_placeholder_replacements,
    build_presentation_filename,
    extract_slide_shapes,
    extract_slides_text,
    get_slide_count,
    inject_content_into_pptx,
    invalidate_image_cache,
    slides_fallback_text,
    slides_to_images,
)
from shared.serializers import build_user_display_name, serialize_comment, serialize_user_brief
from shared.time_utils import utc_now

from .permissions import (
    can_access_presentation,
    can_comment_on_presentation,
    can_return_presentation_for_corrections,
    can_submit_presentation,
    can_validate_presentation,
)
from .service import (
    build_comment_visibility_user_ids,
    build_presentation_response,
    extract_manager_user_ids,
    resolve_next_submission_target,
    resolve_user_from_reference,
)

presentations_bp = Blueprint("presentations", __name__, url_prefix="/api")


@presentations_bp.route("/generate", methods=["POST"])
@jwt_required()
def generate_presentation():
    """
    POST /api/generate
    Body: { "template_filename": "...", "form": { clientName, sector, ... } }
    Returns: { "presentation_id": "...", "slide_count": N }
    """
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    body = request.get_json(force=True)
    tmpl_file = body.get("template_filename")
    form_data = body.get("form", {})

    if not tmpl_file:
        return jsonify({"error": "template_filename is required"}), 400

    tmpl_path = Path(LIBRARY_DIR) / tmpl_file
    if not tmpl_path.exists():
        return jsonify({"error": f"Template not found: {tmpl_file}"}), 404

    try:
        structure = extract_slides_text(tmpl_path)
    except Exception as e:
        return jsonify({"error": f"Cannot read template: {e}"}), 500

    try:
        generated = apply_placeholder_replacements(structure, form_data)
    except Exception as e:
        return jsonify({"error": f"Placeholder replacement error: {e}"}), 500

    pres_id = str(uuid.uuid4())
    out_name = build_presentation_filename(form_data)
    out_file = output_path / f"{pres_id}.pptx"
    try:
        inject_content_into_pptx(tmpl_path, generated, out_file)
    except Exception as e:
        return jsonify({"error": f"PPTX creation failed: {e}"}), 500

    slide_count = get_slide_count(out_file)

    now = utc_now()
    doc = {
        "presentation_id": pres_id,
        "filename":        out_name,
        "path":            str(out_file),
        "template":        tmpl_file,
        "form":            form_data,
        "slide_count":     slide_count,
        "status":          "draft",
        "owner_user_id":   str(current_user["_id"]),
        "owner_name":      build_user_display_name(current_user),
        "current_reviewer_user_id": None,
        "current_reviewer_name": None,
        "current_reviewer_grade": None,
        "created_at":      now,
        "updated_at":      now,
    }
    result = presentations_col.insert_one(doc)

    try:
        slides_to_images(out_file, pres_id)
    except Exception:
        pass

    return jsonify({
        "presentation_id": pres_id,
        "mongo_id":        str(result.inserted_id),
        "filename":        out_name,
        "slide_count":     slide_count,
        "generation_mode": "placeholder-replacement",
    }), 201


@presentations_bp.route("/presentations", methods=["GET"])
@jwt_required()
def list_presentations():
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    user_id = str(current_user["_id"])
    docs = list(
        presentations_col.find(
            {
                "$or": [
                    {"owner_user_id": user_id},
                    {"current_reviewer_user_id": user_id},
                    {"submission_history.submitted_to_user_id": user_id},
                ]
            }
        ).sort("created_at", -1)
    )
    return jsonify({"items": [build_presentation_response(d, current_user) for d in docs]})


@presentations_bp.route("/presentations/<pres_id>", methods=["GET"])
@jwt_required()
def get_presentation(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403
    current_user_id = str(current_user["_id"])
    if any(
        isinstance(notification, dict)
        and str(notification.get("user_id") or "") == current_user_id
        and notification.get("read_at") is None
        for notification in (doc.get("notifications") or [])
    ):
        read_at = utc_now()
        presentations_col.update_one(
            {"presentation_id": pres_id},
            {"$set": {"notifications.$[item].read_at": read_at}},
            array_filters=[{"item.user_id": current_user_id, "item.read_at": None}],
        )
        doc = presentations_col.find_one({"presentation_id": pres_id}) or doc
    return jsonify(build_presentation_response(doc, current_user))


@presentations_bp.route("/presentations/<pres_id>/slides", methods=["GET"])
@jwt_required()
def presentation_slides(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403
    fp = Path(doc["path"])
    if not fp.exists():
        return jsonify({"error": "PPTX file missing on disk"}), 404
    mode = request.args.get("mode", "images")
    if mode == "text":
        try:
            s = slides_fallback_text(fp)
            return jsonify({"slide_count": len(s), "mode": "text", "slides": s})
        except Exception as e:
            return jsonify({"error": str(e)}), 500
    try:
        imgs = slides_to_images(fp, pres_id)
        return jsonify({"slide_count": len(imgs), "mode": "images", "slides": imgs})
    except RuntimeError as e:
        return jsonify({"error": str(e), "fallback_url": f"/api/presentations/{pres_id}/slides?mode=text"}), 503
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@presentations_bp.route("/presentations/<pres_id>/slides/<int:slide_index>/shapes", methods=["GET"])
@jwt_required()
def presentation_slide_shapes(pres_id: str, slide_index: int):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403

    fp = Path(doc["path"])
    if not fp.exists():
        return jsonify({"error": "PPTX file missing on disk"}), 404

    try:
        shapes = extract_slide_shapes(fp, slide_index)
    except IndexError as e:
        return jsonify({"error": str(e)}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    return jsonify({
        "slide_index": slide_index,
        "shapes": shapes,
    })


@presentations_bp.route("/presentations/<pres_id>/slide/<int:slide_index>", methods=["PATCH"])
@jwt_required()
def patch_slide(pres_id: str, slide_index: int):
    """
    PATCH /api/presentations/<id>/slide/<n>
    Body: { "shapes": [{ "shape_id": N, "texts": ["para1", "para2"] }] }
    Saves the PPTX and invalidates the image cache for this presentation.
    """
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403
    fp = Path(doc["path"])
    if not fp.exists():
        return jsonify({"error": "PPTX missing"}), 404

    shapes_patch = {s["shape_id"]: s["texts"] for s in request.get_json(force=True).get("shapes", [])}
    try:
        prs   = Presentation(str(fp))
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.shape_id not in shapes_patch or not shape.has_text_frame:
                continue
            new_texts = shapes_patch[shape.shape_id]
            for i, para in enumerate(shape.text_frame.paragraphs):
                t = new_texts[i] if i < len(new_texts) else ""
                if para.runs:
                    para.runs[0].text = t
                    for run in para.runs[1:]:
                        run.text = ""
        prs.save(str(fp))
    except Exception as e:
        return jsonify({"error": f"Patch failed: {e}"}), 500

    invalidate_image_cache(pres_id)
    presentations_col.update_one(
        {"presentation_id": pres_id},
        {"$set": {"updated_at": utc_now()}},
    )
    return jsonify({"ok": True})


@presentations_bp.route("/presentations/<pres_id>", methods=["PATCH"])
@jwt_required()
def rename_presentation(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403

    payload = request.get_json(force=True) or {}
    raw_filename = str(payload.get("filename") or "").strip()
    if not raw_filename:
        return jsonify({"error": "filename is required"}), 400

    clean_name = raw_filename[:-5] if raw_filename.lower().endswith(".pptx") else raw_filename
    clean_name = re.sub(r"[\\/:*?\"<>|]+", "_", clean_name).strip(" ._")
    if not clean_name:
        return jsonify({"error": "Invalid filename"}), 400

    final_filename = f"{clean_name}.pptx"
    presentations_col.update_one(
        {"presentation_id": pres_id},
        {"$set": {"filename": final_filename, "updated_at": utc_now()}},
    )
    return jsonify({"ok": True, "filename": final_filename})


@presentations_bp.route("/presentations/<pres_id>/comments", methods=["POST"])
@jwt_required()
def add_presentation_comment(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_comment_on_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403

    payload = request.get_json(force=True) or {}
    try:
        slide_index = int(payload.get("slide_index"))
    except (TypeError, ValueError):
        return jsonify({"error": "slide_index is required"}), 400

    comment_text = str(payload.get("text") or "").strip()
    if not comment_text:
        return jsonify({"error": "Comment text is required"}), 400

    timestamp = utc_now()
    comment = {
        "comment_id": str(uuid.uuid4()),
        "slide_index": slide_index,
        "text": comment_text,
        "author_user_id": str(current_user["_id"]),
        "author_name": build_user_display_name(current_user),
        "author_grade": current_user.get("grade"),
        "visible_to_user_ids": build_comment_visibility_user_ids(doc, current_user),
        "created_at": timestamp,
        "updated_at": timestamp,
    }
    notification_recipients = [
        user_id
        for user_id in comment["visible_to_user_ids"]
        if user_id != str(current_user["_id"])
    ]
    notifications = [
        {
            "notification_id": str(uuid.uuid4()),
            "user_id": user_id,
            "presentation_id": pres_id,
            "type": "comment_added",
            "message": f"{build_user_display_name(current_user)} a ajouté un commentaire sur la slide {slide_index + 1}.",
            "slide_index": slide_index,
            "comment_id": comment["comment_id"],
            "created_at": timestamp,
            "read_at": None,
        }
        for user_id in notification_recipients
    ]
    presentations_col.update_one(
        {"presentation_id": pres_id},
        {
            "$push": {
                "slide_comments": comment,
                "notifications": {"$each": notifications},
            },
            "$set": {"updated_at": timestamp},
        },
    )
    return jsonify({"ok": True, "comment": serialize_comment(comment)}), 201


@presentations_bp.route("/presentations/<pres_id>/return", methods=["POST"])
@jwt_required()
def return_presentation_for_corrections(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_return_presentation_for_corrections(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403

    owner_id = str(doc.get("owner_user_id") or "")
    if not owner_id:
        return jsonify({"error": "Presentation owner not found"}), 400
    owner_user = User.collection.find_one({"_id": ObjectId(owner_id)}) if ObjectId.is_valid(owner_id) else None

    timestamp = utc_now()
    return_notification = {
        "notification_id": str(uuid.uuid4()),
        "user_id": owner_id,
        "presentation_id": pres_id,
        "type": "presentation_returned",
        "message": f"{build_user_display_name(current_user)} a renvoyé la présentation pour correction.",
        "created_at": timestamp,
        "read_at": None,
    }
    presentations_col.update_one(
        {"presentation_id": pres_id},
        {
            "$set": {
                "status": "changes_requested",
                "updated_at": timestamp,
                "current_reviewer_user_id": owner_id,
                "current_reviewer_name": doc.get("owner_name"),
                "current_reviewer_grade": owner_user.get("grade") if owner_user else None,
                "returned_by_user_id": str(current_user["_id"]),
                "returned_by_name": build_user_display_name(current_user),
            },
            "$push": {
                "notifications": return_notification,
                "submission_history": {
                    "submitted_at": timestamp,
                    "submitted_by_user_id": str(current_user["_id"]),
                    "submitted_by_name": build_user_display_name(current_user),
                    "submitted_to_user_id": owner_id,
                    "submitted_to_name": doc.get("owner_name"),
                    "submitted_to_role": "assistant",
                    "status": "changes_requested",
                    "action": "returned_for_corrections",
                }
            },
        },
    )
    return jsonify({"ok": True, "status": "changes_requested"})


@presentations_bp.route("/presentations/<pres_id>/submit", methods=["POST"])
@jwt_required()
def submit_presentation(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403
    if not can_submit_presentation(doc, current_user):
        return jsonify({"error": "You cannot submit this presentation at this stage"}), 403

    next_target = resolve_next_submission_target(current_user, doc)
    if not next_target:
        return jsonify({"error": "This presentation cannot be submitted further"}), 400

    target_role, target_reference, allowed_grades, next_status, _label = next_target
    target_user = resolve_user_from_reference(target_reference, allowed_grades=allowed_grades)
    if not target_user:
        return jsonify({"error": f"Unable to find the selected {target_role} in active users"}), 400
    if str(target_user["_id"]) == str(current_user["_id"]):
        return jsonify({"error": "You cannot submit a presentation to yourself"}), 400

    timestamp = utc_now()
    submit_notification = {
        "notification_id": str(uuid.uuid4()),
        "user_id": str(target_user["_id"]),
        "presentation_id": pres_id,
        "type": "presentation_submitted",
        "message": f"{build_user_display_name(current_user)} vous a soumis une présentation.",
        "created_at": timestamp,
        "read_at": None,
    }
    presentations_col.update_one(
        {"presentation_id": pres_id},
        {
            "$set": {
                "status": next_status,
                "submitted_at": timestamp,
                "updated_at": timestamp,
                "current_reviewer_user_id": str(target_user["_id"]),
                "current_reviewer_name": build_user_display_name(target_user),
                "current_reviewer_grade": target_user.get("grade"),
                "submitted_by_user_id": str(current_user["_id"]),
                "submitted_by_name": build_user_display_name(current_user),
            },
            "$push": {
                "notifications": submit_notification,
                "submission_history": {
                    "submitted_at": timestamp,
                    "submitted_by_user_id": str(current_user["_id"]),
                    "submitted_by_name": build_user_display_name(current_user),
                    "submitted_to_user_id": str(target_user["_id"]),
                    "submitted_to_name": build_user_display_name(target_user),
                    "submitted_to_role": target_role,
                    "status": next_status,
                }
            },
        },
    )
    return jsonify({
        "ok": True,
        "status": next_status,
        "submitted_to": serialize_user_brief(target_user),
    })


@presentations_bp.route("/presentations/<pres_id>/validate", methods=["POST"])
@jwt_required()
def validate_presentation(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_validate_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403

    timestamp = utc_now()
    owner_id = str(doc.get("owner_user_id") or "").strip()
    notifications = []
    recipient_ids = {owner_id, *extract_manager_user_ids(doc)}
    recipient_ids.discard(str(current_user["_id"]))
    for user_id in recipient_ids:
        if not user_id:
            continue
        notifications.append({
            "notification_id": str(uuid.uuid4()),
            "user_id": user_id,
            "presentation_id": pres_id,
            "type": "presentation_validated",
            "message": f"{build_user_display_name(current_user)} a validé la présentation.",
            "created_at": timestamp,
            "read_at": None,
        })

    presentations_col.update_one(
        {"presentation_id": pres_id},
        {
            "$set": {
                "status": "completed",
                "validated_at": timestamp,
                "validated_by_user_id": str(current_user["_id"]),
                "validated_by_name": build_user_display_name(current_user),
                "updated_at": timestamp,
                "current_reviewer_user_id": None,
                "current_reviewer_name": None,
                "current_reviewer_grade": None,
            },
            "$push": {
                "notifications": {"$each": notifications},
                "submission_history": {
                    "submitted_at": timestamp,
                    "submitted_by_user_id": str(current_user["_id"]),
                    "submitted_by_name": build_user_display_name(current_user),
                    "submitted_to_user_id": None,
                    "submitted_to_name": None,
                    "submitted_to_role": "completed",
                    "status": "completed",
                    "action": "validated",
                },
            },
        },
    )
    return jsonify({"ok": True, "status": "completed"})


@presentations_bp.route("/presentations/<pres_id>/download", methods=["GET"])
@jwt_required()
def download_presentation(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403
    fp = Path(doc["path"])
    if not fp.exists():
        return jsonify({"error": "File not found"}), 404
    name = doc.get("filename") or f"Propale_{doc['form'].get('clientName', 'client')}.pptx"
    return send_file(
        fp,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=name,
    )
