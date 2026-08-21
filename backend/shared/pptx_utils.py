import base64
import glob
import re
import shutil
import subprocess
from pathlib import Path

from bson import ObjectId
from pptx import Presentation

from extensions import LIBREOFFICE_PATH, cache_path, clients_col

from .text_utils import normalize_placeholder_key, slugify_filename_part
from .time_utils import utc_now


def find_soffice() -> str | None:
    if LIBREOFFICE_PATH and Path(LIBREOFFICE_PATH).exists():
        return LIBREOFFICE_PATH
    found = shutil.which("soffice")
    if found:
        return found
    for c in [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        *glob.glob(r"C:\Program Files\LibreOffice*\program\soffice.exe"),
        "/usr/bin/soffice",
        "/usr/local/bin/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        *glob.glob("/opt/libreoffice*/program/soffice"),
    ]:
        if c and Path(c).exists():
            return c
    return None


def get_slide_count(file_path: Path) -> int:
    try:
        return len(Presentation(str(file_path)).slides)
    except Exception:
        return 0


def extract_slides_text(file_path: Path) -> list[dict]:
    """
    Extract every text shape from a PPTX.
    Returns a list usable both as Claude input and as injection blueprint.
    """
    prs = Presentation(str(file_path))
    result = []
    for slide_idx, slide in enumerate(prs.slides):
        shapes_data = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            texts = [p.text for p in shape.text_frame.paragraphs]
            if not any(t.strip() for t in texts):
                continue
            ph_idx = None
            try:
                if shape.is_placeholder:
                    ph_idx = shape.placeholder_format.idx
            except Exception:
                ph_idx = None
            shapes_data.append({
                "shape_id":        shape.shape_id,
                "shape_name":      shape.name,
                "placeholder_idx": ph_idx,
                "texts":           texts,
            })
        result.append({"slide_index": slide_idx, "shapes": shapes_data})
    return result


def extract_slide_shapes(file_path: Path, slide_index: int) -> list[dict]:
    prs = Presentation(str(file_path))
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError("Slide index out of range")

    shapes_data = []
    slide = prs.slides[slide_index]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        texts = [p.text for p in shape.text_frame.paragraphs]
        if not any(t.strip() for t in texts):
            continue
        shapes_data.append({
            "shape_id": shape.shape_id,
            "shape_name": shape.name,
            "texts": texts,
        })
    return shapes_data


def build_presentation_filename(form_data: dict) -> str:
    mission = slugify_filename_part((form_data or {}).get("missionType"), "mission")
    client = slugify_filename_part((form_data or {}).get("clientName"), "client")
    return f"Propale_{mission}_{client}.pptx"


def build_placeholder_context(form_data: dict) -> dict[str, str]:
    context: dict[str, str] = {}

    for key, value in (form_data or {}).items():
        if value is None:
            continue
        if isinstance(value, list):
            string_value = ", ".join(str(item).strip() for item in value if str(item).strip())
        if isinstance(value, (str, int, float)):
            string_value = str(value).strip()
        elif not isinstance(value, list):
            string_value = str(value)
        if string_value == "":
            continue
        normalized_key = normalize_placeholder_key(key)
        if normalized_key:
            context[normalized_key] = string_value

    aliases = {
        "client": context.get("clientname", ""),
        "client_name": context.get("clientname", ""),
        "company_name": context.get("clientname", ""),
        "nom_client": context.get("clientname", ""),
        "secteur": context.get("sector", ""),
        "pays": context.get("country", ""),
        "mission": context.get("missiontype", ""),
        "mission_type": context.get("missiontype", ""),
        "normes": context.get("standards", ""),
        "contexte": context.get("context", ""),
        "objectifs": context.get("objectives", ""),
        "associe": context.get("partner", ""),
        "manager_name": context.get("manager", ""),
        "honoraires": context.get("fees", ""),
        "date_limite": context.get("deadline", ""),
        "duree": context.get("duration", ""),
    }
    for alias, value in aliases.items():
        if value:
            context[alias] = value

    client_id = form_data.get("clientId")
    if client_id:
        try:
            client_doc = clients_col.find_one({"_id": ObjectId(client_id)})
        except Exception:
            client_doc = None
        if client_doc:
            client_fields = {
                "client_company_name": client_doc.get("company_name", ""),
                "client_sector": client_doc.get("sector", ""),
                "client_country": client_doc.get("country", ""),
                "client_civility": client_doc.get("civility", ""),
                "client_responsable_name": client_doc.get("responsable_name", ""),
                "client_responsable_function": client_doc.get("responsable_function", ""),
                "client_legal_form": client_doc.get("legal_form", ""),
                "client_rccm": client_doc.get("RCCM", ""),
                "client_address": client_doc.get("address", ""),
                "client_creation_date": client_doc.get("creation_date", ""),
                "responsable_name": client_doc.get("responsable_name", ""),
                "responsable_function": client_doc.get("responsable_function", ""),
                "legal_form": client_doc.get("legal_form", ""),
                "rccm": client_doc.get("RCCM", ""),
                "address": client_doc.get("address", ""),
            }
            for key, value in client_fields.items():
                if value:
                    context[key] = str(value).strip()

    now = utc_now()
    context.setdefault("generated_date", now.strftime("%Y-%m-%d"))
    context.setdefault("generated_at", now.strftime("%Y-%m-%d %H:%M"))
    context.setdefault("current_year", str(now.year))
    return context


PLACEHOLDER_PATTERNS = [
    re.compile(r"\{\{\s*([a-zA-Z0-9_.-]+)\s*\}\}"),
    re.compile(r"\[\[\s*([a-zA-Z0-9_.-]+)\s*\]\]"),
    re.compile(r"<<\s*([a-zA-Z0-9_.-]+)\s*>>"),
    re.compile(r"%\s*([a-zA-Z0-9_.-]+)\s*%"),
]


def replace_placeholders_in_text(text: str, context: dict[str, str]) -> str:
    replaced = text
    for pattern in PLACEHOLDER_PATTERNS:
        replaced = pattern.sub(
            lambda match: context.get(normalize_placeholder_key(match.group(1)), match.group(0)),
            replaced,
        )
    return replaced


def apply_placeholder_replacements(slides_structure: list[dict], form_data: dict) -> list[dict]:
    context = build_placeholder_context(form_data)
    generated_content = []

    for slide_data in slides_structure:
        generated_shapes = []
        for shape_data in slide_data["shapes"]:
            generated_shapes.append({
                **shape_data,
                "texts": [
                    replace_placeholders_in_text(text, context)
                    for text in shape_data["texts"]
                ],
            })
        generated_content.append({
            "slide_index": slide_data["slide_index"],
            "shapes": generated_shapes,
        })

    return generated_content


def inject_content_into_pptx(template_path: Path, content: list[dict], out_file: Path):
    """
    Copy template and replace paragraph texts while preserving all formatting.
    `content` has the same shape as extract_slides_text() output.
    """
    prs = Presentation(str(template_path))

    lookup: dict[tuple, list[str]] = {}
    for slide_data in content:
        for shape_data in slide_data["shapes"]:
            lookup[(slide_data["slide_index"], shape_data["shape_id"])] = shape_data["texts"]

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            new_texts = lookup.get((slide_idx, shape.shape_id))
            if new_texts is None or not shape.has_text_frame:
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                new_text = new_texts[para_idx] if para_idx < len(new_texts) else ""
                if para.runs:
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ""

    prs.save(str(out_file))


def ensure_pdf(file_path: Path, stem: str | None = None) -> Path:
    stem = stem or file_path.stem
    pdf_path = cache_path / f"{stem}.pdf"
    if pdf_path.exists():
        return pdf_path
    soffice = find_soffice()
    if not soffice:
        raise RuntimeError("LibreOffice not found. Install it or set LIBREOFFICE_PATH.")
    tmp = cache_path / f"{stem}.pptx"
    if str(tmp) != str(file_path):
        shutil.copy2(str(file_path), str(tmp))
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(cache_path), str(tmp)],
        check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
    )
    return pdf_path


def pdf_to_images(pdf_path: Path, stem: str) -> list[str]:
    prefix = cache_path / f"{stem}_slide"
    pdftoppm = shutil.which("pdftoppm")
    if pdftoppm:
        subprocess.run(
            [pdftoppm, "-png", "-r", "120", str(pdf_path), str(prefix)],
            check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        return [base64.b64encode(p.read_bytes()).decode()
                for p in sorted(cache_path.glob(f"{stem}_slide-*.png"))]
    try:
        import fitz
        doc = fitz.open(str(pdf_path))
        return [base64.b64encode(page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5)).tobytes("png")).decode()
                for page in doc]
    except ImportError:
        pass
    raise RuntimeError("Install poppler-utils (pdftoppm) or pymupdf.")


def slides_to_images(pptx_path: Path, stem: str) -> list[str]:
    existing = sorted(cache_path.glob(f"{stem}_slide-*.png"))
    if existing:
        return [base64.b64encode(p.read_bytes()).decode() for p in existing]
    return pdf_to_images(ensure_pdf(pptx_path, stem=stem), stem)


def slides_fallback_text(file_path: Path) -> list[dict]:
    prs = Presentation(str(file_path))
    out = []
    for i, slide in enumerate(prs.slides):
        title, texts, img = "", [], None
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    ph_idx = None
                    try:
                        if shape.is_placeholder:
                            ph_idx = shape.placeholder_format.idx
                    except Exception:
                        ph_idx = None
                    if not title and ph_idx == 0:
                        title = t
                    else:
                        texts.append(t)
            if img is None and shape.shape_type == 13:
                try:
                    img = base64.b64encode(shape.image.blob).decode()
                except Exception:
                    pass
        out.append({"index": i + 1, "title": title, "text": "\n".join(texts), "image": img})
    return out


def invalidate_image_cache(stem: str):
    for f in cache_path.glob(f"{stem}_slide-*.png"):
        f.unlink(missing_ok=True)
    (cache_path / f"{stem}.pdf").unlink(missing_ok=True)
