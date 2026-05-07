import io
import os
import re
import uuid
import base64
import glob
import shutil
import subprocess
import unicodedata
from datetime import UTC, datetime, timedelta
from pathlib import Path
from pptx.util import Inches
from flask import Flask, jsonify, request, send_file
from flask_cors import CORS
from pymongo import MongoClient, UpdateOne
from pptx import Presentation
from dotenv import load_dotenv
from flask_jwt_extended import JWTManager, jwt_required, create_access_token, get_jwt_identity
from config import Config
from models import User, Proposal, Element
import time
import bcrypt
from bson import ObjectId

load_dotenv()

app = Flask(__name__)
app.config.from_object(Config)
CORS(app)
jwt = JWTManager(app)

MONGO_URI = os.getenv("MONGO_URI", Config.MONGO_URI)
LIBRARY_DIR = os.getenv("LIBRARY_DIR", str(Path(__file__).resolve().parents[1] / "Propale_library"))
CACHE_DIR = os.getenv("CACHE_DIR", str(Path(__file__).resolve().parents[1] / "Propale_cache"))
OUTPUT_DIR = os.getenv("OUTPUT_DIR", str(Path(__file__).resolve().parents[1] / "Propale_output"))
LIBREOFFICE_PATH  = os.getenv("LIBREOFFICE_PATH",  "")

mongo_client      = MongoClient(MONGO_URI)
db                = mongo_client.get_default_database()
templates_col     = db["templates"]
presentations_col = db["presentations"]
clients_col = db["clients"]

cache_path  = Path(CACHE_DIR)
output_path = Path(OUTPUT_DIR)
cache_path.mkdir(parents=True, exist_ok=True)
output_path.mkdir(parents=True, exist_ok=True)


# Ensure uploads directory exists
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)


def serialize_client(doc: dict) -> dict:
    if doc and "_id" in doc:
        doc["_id"] = str(doc["_id"])
    return doc


def utc_now() -> datetime:
    return datetime.now(UTC)


def ensure_utc_datetime(value):
    if isinstance(value, datetime):
        if value.tzinfo is None:
            return value.replace(tzinfo=UTC)
        return value.astimezone(UTC)
    return None
 
 
@app.route("/api/clients", methods=["GET"])
def list_clients():
    """
    GET /api/clients?page=1&limit=25&q=search
    Recherche sur company_name, sector, country, responsable_name.
    """
    page  = max(1, int(request.args.get("page", 1)))
    limit = min(100, int(request.args.get("limit", 25)))
    q     = request.args.get("q", "").strip()
    skip  = (page - 1) * limit
 
    query = {}
    if q:
        query["$or"] = [
            {"company_name":       {"$regex": q, "$options": "i"}},
            {"sector":             {"$regex": q, "$options": "i"}},
            {"country":            {"$regex": q, "$options": "i"}},
            {"responsable_name":   {"$regex": q, "$options": "i"}},
        ]
 
    total = clients_col.count_documents(query)
    docs  = list(
        clients_col.find(query)
        .sort("company_name", 1)
        .skip(skip)
        .limit(limit)
    )
 
    return jsonify({
        "items":      [serialize_client(d) for d in docs],
        "total":      total,
        "page":       page,
        "limit":      limit,
        "totalPages": max(1, (total + limit - 1) // limit),
    })
 
 
@app.route("/api/clients", methods=["POST"])
def create_client():
    data = request.get_json(force=True)
 
    if not data.get("company_name", "").strip():
        return jsonify({"error": "Le nom du client est requis"}), 400
 
    now = utc_now()
    doc = {
        "company_name":        data.get("company_name", "").strip(),
        "sector":              data.get("sector", ""),
        "country":             data.get("country", ""),
        "civility":            data.get("civility", ""),
        "responsable_name":    data.get("responsable_name", ""),
        "responsable_function":data.get("responsable_function", ""),
        "legal_form":          data.get("legal_form", ""),
        "RCCM":                data.get("RCCM", ""),
        "address":             data.get("address", ""),
        "creation_date":       data.get("creation_date", ""),
        "created_at":          now,
        "updated_at":          now,
    }
    result = clients_col.insert_one(doc)
    doc["_id"] = str(result.inserted_id)
    return jsonify(doc), 201
 
 
@app.route("/api/clients/<client_id>", methods=["GET"])
def get_client(client_id: str):
    from bson import ObjectId
    try:
        doc = clients_col.find_one({"_id": ObjectId(client_id)})
    except Exception:
        return jsonify({"error": "Invalid ID"}), 400
    if not doc:
        return jsonify({"error": "Client not found"}), 404
    return jsonify(serialize_client(doc))
 
 
@app.route("/api/clients/<client_id>", methods=["PUT"])
def update_client(client_id: str):
    from bson import ObjectId
    data = request.get_json(force=True)
 
    if not data.get("company_name", "").strip():
        return jsonify({"error": "Le nom du client est requis"}), 400
 
    updates = {
        "company_name":        data.get("company_name", "").strip(),
        "sector":              data.get("sector", ""),
        "country":             data.get("country", ""),
        "civility":            data.get("civility", ""),
        "responsable_name":    data.get("responsable_name", ""),
        "responsable_function":data.get("responsable_function", ""),
        "legal_form":          data.get("legal_form", ""),
        "RCCM":                data.get("RCCM", ""),
        "address":             data.get("address", ""),
        "creation_date":       data.get("creation_date", ""),
        "updated_at":          utc_now(),
    }
 
    try:
        result = clients_col.update_one({"_id": ObjectId(client_id)}, {"$set": updates})
    except Exception:
        return jsonify({"error": "Invalid ID"}), 400
 
    if result.matched_count == 0:
        return jsonify({"error": "Client not found"}), 404
 
    doc = clients_col.find_one({"_id": ObjectId(client_id)})
    return jsonify(serialize_client(doc))
 
 
@app.route("/api/clients/<client_id>", methods=["DELETE"])
def delete_client(client_id: str):
    from bson import ObjectId
    try:
        result = clients_col.delete_one({"_id": ObjectId(client_id)})
    except Exception:
        return jsonify({"error": "Invalid ID"}), 400
    if result.deleted_count == 0:
        return jsonify({"error": "Client not found"}), 404
    return jsonify({"ok": True})
 
 
@app.route("/api/clients/meta/sectors", methods=["GET"])
def client_sectors():
    raw = clients_col.distinct("sector")
    sectors = sorted([s for s in raw if s and s.strip()])
    return jsonify({"items": sectors})
 
 
@app.route("/api/clients/meta/countries", methods=["GET"])
def client_countries():
    raw = clients_col.distinct("country")
    countries = sorted([c for c in raw if c and c.strip()])
    return jsonify({"items": countries})


def validate_password_strength(password):
    if len(password) < 8:
        return "Password must contain at least 8 characters"
    if not re.search(r"[A-Z]", password):
        return "Password must contain at least one uppercase letter"
    if not re.search(r"[a-z]", password):
        return "Password must contain at least one lowercase letter"
    if not re.search(r"\d", password):
        return "Password must contain at least one number"
    if not re.search(r"[!@#$%^&*()_+\-=\[\]{};':\"\\|,.<>\/?]", password):
        return "Password must contain at least one special character"
    return None


def normalize_email(email):
    if email is None:
        return None
    return str(email).strip().lower()


ASSISTANT_SUBMISSION_GRADES = {"Assistant", "Senior"}
MANAGER_SUBMISSION_GRADES = {"Assistant manager", "Manager", "Senior manager"}
ASSOCIATE_GRADES = {"Associé"}
FORM_MANAGER_GRADES = {"Manager", "Senior manager"}
MANAGER_REVIEWER_GRADES = FORM_MANAGER_GRADES
COMMENTER_GRADES = ASSISTANT_SUBMISSION_GRADES | MANAGER_SUBMISSION_GRADES | ASSOCIATE_GRADES


def normalize_person_label(value):
    if value is None:
        return ""
    text = unicodedata.normalize("NFKD", str(value)).encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^a-z0-9]+", " ", text.lower()).strip()
    return re.sub(r"\s+", " ", text)


def slugify_filename_part(value: str | None, fallback: str) -> str:
    if value is None:
        return fallback
    text = unicodedata.normalize("NFKD", str(value)).encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^A-Za-z0-9]+", "_", text).strip("_")
    return text or fallback


def build_presentation_filename(form_data: dict) -> str:
    mission = slugify_filename_part((form_data or {}).get("missionType"), "mission")
    client = slugify_filename_part((form_data or {}).get("clientName"), "client")
    return f"Propale_{mission}_{client}.pptx"


def build_user_display_name(user: dict | None) -> str | None:
    if not user:
        return None
    if user.get("name"):
        return str(user["name"]).strip()
    parts = [str(user.get("first_name") or "").strip(), str(user.get("last_name") or "").strip()]
    full_name = " ".join(part for part in parts if part)
    return full_name or user.get("email")


def serialize_timestamp(value) -> str | None:
    if isinstance(value, datetime):
        return value.isoformat()
    return value


def serialize_comment(comment: dict) -> dict:
    payload = dict(comment or {})
    payload["created_at"] = serialize_timestamp(payload.get("created_at"))
    payload["updated_at"] = serialize_timestamp(payload.get("updated_at"))
    return payload


def serialize_notification(notification: dict) -> dict:
    payload = dict(notification or {})
    payload["created_at"] = serialize_timestamp(payload.get("created_at"))
    payload["read_at"] = serialize_timestamp(payload.get("read_at"))
    return payload


def serialize_user_brief(user: dict) -> dict:
    return {
        "id": str(user["_id"]),
        "name": build_user_display_name(user),
        "email": user.get("email"),
        "grade": user.get("grade"),
        "code_categorie": user.get("code_categorie"),
    }


def get_current_user_doc():
    user_id = get_jwt_identity()
    if not user_id:
        return None
    try:
        return User.find_by_id(user_id)
    except Exception:
        return None


def resolve_user_from_reference(reference: str | None, allowed_grades: set[str] | None = None):
    if not reference:
        return None

    normalized_ref = normalize_person_label(reference)
    if not normalized_ref:
        return None

    for user in User.collection.find({"is_active": True}):
        if allowed_grades and user.get("grade") not in allowed_grades:
            continue

        candidates = {
            normalize_person_label(user.get("email")),
            normalize_person_label(user.get("name")),
            normalize_person_label(
                f"{str(user.get('first_name') or '').strip()} {str(user.get('last_name') or '').strip()}"
            ),
            normalize_person_label(
                f"{str(user.get('last_name') or '').strip()} {str(user.get('first_name') or '').strip()}"
            ),
        }
        if normalized_ref in candidates:
            return user
    return None


def can_access_presentation(doc: dict, user: dict) -> bool:
    user_id = str(user["_id"])
    history = doc.get("submission_history") or []
    has_been_reviewer = any(
        str(entry.get("submitted_to_user_id") or "") == user_id
        for entry in history
        if isinstance(entry, dict)
    )
    return (
        str(doc.get("owner_user_id") or "") == user_id
        or str(doc.get("current_reviewer_user_id") or "") == user_id
        or has_been_reviewer
    )


def can_submit_presentation(doc: dict, user: dict) -> bool:
    if not can_access_presentation(doc, user):
        return False

    user_id = str(user["_id"])
    user_grade = user.get("grade")
    current_reviewer_id = str(doc.get("current_reviewer_user_id") or "")
    owner_id = str(doc.get("owner_user_id") or "")

    if user_grade in ASSOCIATE_GRADES:
        return False

    if user_grade in ASSISTANT_SUBMISSION_GRADES:
        return owner_id == user_id

    if user_grade in MANAGER_SUBMISSION_GRADES:
        return True

    if current_reviewer_id:
        return current_reviewer_id == user_id

    return owner_id == user_id


def can_comment_on_presentation(doc: dict, user: dict) -> bool:
    return can_access_presentation(doc, user) and user.get("grade") in COMMENTER_GRADES


def can_return_presentation_for_corrections(doc: dict, user: dict) -> bool:
    return (
        can_access_presentation(doc, user)
        and str(doc.get("current_reviewer_user_id") or "") == str(user["_id"])
        and str(doc.get("owner_user_id") or "") != str(user["_id"])
        and user.get("grade") in MANAGER_SUBMISSION_GRADES
    )


def can_validate_presentation(doc: dict, user: dict) -> bool:
    return (
        can_access_presentation(doc, user)
        and str(doc.get("current_reviewer_user_id") or "") == str(user["_id"])
        and user.get("grade") in ASSOCIATE_GRADES
        and str(doc.get("status") or "") == "submitted_to_associate"
    )


def extract_manager_user_ids(doc: dict) -> set[str]:
    manager_ids: set[str] = set()
    history = doc.get("submission_history") or []
    for entry in history:
        if not isinstance(entry, dict):
            continue
        if entry.get("submitted_to_role") == "manager":
            user_id = str(entry.get("submitted_to_user_id") or "").strip()
            if user_id:
                manager_ids.add(user_id)
        if entry.get("status") == "submitted_to_associate":
            user_id = str(entry.get("submitted_by_user_id") or "").strip()
            if user_id:
                manager_ids.add(user_id)
        if entry.get("action") == "returned_for_corrections":
            user_id = str(entry.get("submitted_by_user_id") or "").strip()
            if user_id:
                manager_ids.add(user_id)
    current_reviewer_id = str(doc.get("current_reviewer_user_id") or "").strip()
    if current_reviewer_id:
        current_grade = str(doc.get("current_reviewer_grade") or "")
        if current_grade in MANAGER_SUBMISSION_GRADES:
            manager_ids.add(current_reviewer_id)
    return manager_ids


def build_comment_visibility_user_ids(doc: dict, author: dict) -> list[str]:
    owner_id = str(doc.get("owner_user_id") or "").strip()
    manager_ids = extract_manager_user_ids(doc)
    author_id = str(author.get("_id") or "").strip()

    visible_ids = {user_id for user_id in [owner_id, author_id] if user_id}
    visible_ids.update(manager_ids)
    return sorted(visible_ids)


def get_visible_slide_comments(doc: dict, viewer: dict) -> list[dict]:
    viewer_id = str(viewer["_id"])
    comments = []
    for comment in (doc.get("slide_comments") or []):
        if not isinstance(comment, dict):
            continue
        visible_to = {str(user_id).strip() for user_id in (comment.get("visible_to_user_ids") or []) if str(user_id).strip()}
        author_id = str(comment.get("author_user_id") or "").strip()
        if viewer_id == author_id or not visible_to or viewer_id in visible_to:
            comments.append(serialize_comment(comment))
    return sorted(comments, key=lambda item: item.get("created_at") or "")


def build_notifications_for_user(doc: dict, user_id: str) -> list[dict]:
    notifications = []
    for notification in (doc.get("notifications") or []):
        if not isinstance(notification, dict):
            continue
        if str(notification.get("user_id") or "") != user_id:
            continue
        payload = serialize_notification(notification)
        payload.setdefault("presentation_id", doc.get("presentation_id"))
        payload.setdefault("filename", doc.get("filename"))
        notifications.append(payload)
    return notifications


def build_presentation_search_text(doc: dict) -> str:
    form = doc.get("form") or {}
    fields = [
        doc.get("filename"),
        doc.get("owner_name"),
        doc.get("template"),
        form.get("missionType"),
        form.get("clientName"),
        form.get("manager"),
        form.get("partner"),
        form.get("sector"),
    ]
    values = []
    for field in fields:
        if isinstance(field, list):
            values.extend(str(item).strip() for item in field if str(item).strip())
        elif field:
            values.append(str(field).strip())
    return " ".join(values).lower()


def resolve_next_submission_target(actor: dict, presentation_doc: dict):
    grade = actor.get("grade")
    form = presentation_doc.get("form") or {}
    partner_reference = form.get("partner")
    if isinstance(partner_reference, list):
        partner_reference = next(
            (str(item).strip() for item in partner_reference if str(item).strip()),
            None,
        )

    if grade in ASSISTANT_SUBMISSION_GRADES:
        return (
            "manager",
            form.get("manager"),
            MANAGER_REVIEWER_GRADES,
            "submitted_to_manager",
            "Soumettre au manager",
        )
    if grade in MANAGER_SUBMISSION_GRADES:
        return (
            "associate",
            partner_reference,
            ASSOCIATE_GRADES,
            "submitted_to_associate",
            "Soumettre à l'associé",
        )
    if grade in ASSOCIATE_GRADES:
        return None
    return None


def build_presentation_response(doc: dict, viewer: dict) -> dict:
    payload = serialize(doc.copy())
    owner_id = str(doc.get("owner_user_id") or "")
    viewer_id = str(viewer["_id"])
    history = doc.get("submission_history") or []
    has_been_reviewer = any(
        str(entry.get("submitted_to_user_id") or "") == viewer_id
        for entry in history
        if isinstance(entry, dict)
    )
    submission_allowed = can_submit_presentation(doc, viewer)
    next_target = resolve_next_submission_target(viewer, doc) if submission_allowed else None
    target_user = None
    if next_target:
        target_user = resolve_user_from_reference(next_target[1], allowed_grades=next_target[2])
    comments = sorted(
        get_visible_slide_comments(doc, viewer),
        key=lambda item: item.get("created_at") or "",
    )
    payload.update({
        "is_owner": owner_id == viewer_id,
        "is_reviewer": str(doc.get("current_reviewer_user_id") or "") == viewer_id,
        "has_been_reviewer": has_been_reviewer,
        "can_submit": target_user is not None and submission_allowed,
        "can_validate": can_validate_presentation(doc, viewer),
        "can_comment": can_comment_on_presentation(doc, viewer),
        "can_return_for_corrections": can_return_presentation_for_corrections(doc, viewer),
        "submit_action_label": next_target[4] if next_target else None,
        "validate_action_label": "Valider la présentation",
        "return_action_label": "Renvoyer à l'assistant pour correction",
        "current_reviewer_name": doc.get("current_reviewer_name"),
        "next_reviewer_name": build_user_display_name(target_user),
        "owner_name": doc.get("owner_name"),
        "slide_comments": comments,
    })
    return payload


@app.route("/api/auth/register", methods=["POST"])
def register():
    data = request.get_json()

    if not data or not data.get('email') or not data.get('password'):
        return jsonify({"error": "Email and password are required"}), 400

    email = normalize_email(data['email'])
    password = data['password']
    name = data.get('name')

    password_error = validate_password_strength(password)
    if password_error:
        return jsonify({"error": password_error}), 400

    # Check if user already exists
    if User.find_by_email(email):
        return jsonify({"error": "User already exists"}), 409

    try:
        user_id = User.create_user(email, password, name)
        return jsonify({
            "message": "User created successfully",
            "user_id": user_id
        }), 201
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/auth/login", methods=["POST"])
def login():
    data = request.get_json()

    if not data or not data.get('email') or not data.get('password'):
        return jsonify({"error": "Email and password are required"}), 400

    email = normalize_email(data['email'])
    password = data['password']

    user = User.find_by_email(email)
    if not user or not user.get("is_active", True) or not User.verify_password(user['password'], password):
        return jsonify({"error": "Invalid credentials"}), 401

    now = utc_now()
    User.collection.update_one(
        {"_id": user["_id"]},
        {"$set": {"last_seen_at": now, "last_login_at": now}},
    )
    access_token = create_access_token(identity=str(user['_id']))
    return jsonify({
        "access_token": access_token,
        "user": {
            "id": str(user['_id']),
            "email": user['email'],
            "name": user.get('name'),
            "first_name": user.get('first_name'),
            "last_name": user.get('last_name'),
            "code_categorie": user.get('code_categorie'),
            "grade": user.get('grade'),
            "department": user.get('department'),
        }
    }), 200


@app.route("/api/users/reviewers", methods=["GET"])
def list_reviewers():
    users = list(User.collection.find({"is_active": True}))
    managers = [
        serialize_user_brief(user)
        for user in users
        if user.get("grade") in MANAGER_REVIEWER_GRADES
    ]
    associates = [
        serialize_user_brief(user)
        for user in users
        if user.get("grade") in ASSOCIATE_GRADES
    ]

    managers.sort(key=lambda item: (item.get("name") or "", item.get("email") or ""))
    associates.sort(key=lambda item: (item.get("name") or "", item.get("email") or ""))
    return jsonify({"managers": managers, "associates": associates})


@app.route("/api/session/heartbeat", methods=["POST"])
@jwt_required()
def session_heartbeat():
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    User.collection.update_one(
        {"_id": current_user["_id"]},
        {"$set": {"last_seen_at": utc_now()}},
    )
    return jsonify({"ok": True})


@app.route("/api/dashboard/summary", methods=["GET"])
@jwt_required()
def dashboard_summary():
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    active_since = utc_now() - timedelta(minutes=5)
    active_count = User.collection.count_documents({
        "is_active": True,
        "last_seen_at": {"$gte": active_since},
    })

    query = str(request.args.get("q") or "").strip().lower()
    docs = list(
        presentations_col.find({}, {
            "_id": 0,
            "presentation_id": 1,
            "filename": 1,
            "form": 1,
            "status": 1,
            "created_at": 1,
            "updated_at": 1,
            "owner_name": 1,
            "owner_user_id": 1,
            "current_reviewer_user_id": 1,
            "submission_history": 1,
            "template": 1,
        }).sort("created_at", -1)
    )

    visible_docs = [doc for doc in docs if can_access_presentation(doc, current_user)]
    if query:
        visible_docs = [doc for doc in visible_docs if query in build_presentation_search_text(doc)]
    limit = 20 if query else 3
    recent_docs = visible_docs[:limit]

    items = []
    for doc in recent_docs:
        created_at = ensure_utc_datetime(doc.get("created_at"))
        updated_at = ensure_utc_datetime(doc.get("updated_at"))
        items.append({
            "presentation_id": doc.get("presentation_id"),
            "title": (doc.get("filename") or "Présentation").replace(".pptx", ""),
            "mission_type": (doc.get("form") or {}).get("missionType"),
            "client_name": (doc.get("form") or {}).get("clientName"),
            "status": doc.get("status"),
            "owner_name": doc.get("owner_name"),
            "created_at": created_at.isoformat() if created_at else None,
            "updated_at": updated_at.isoformat() if updated_at else None,
        })

    return jsonify({
        "active_users_count": active_count,
        "recent_presentations": items,
    })


@app.route("/api/notifications", methods=["GET"])
@jwt_required()
def list_notifications():
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    user_id = str(current_user["_id"])
    docs = list(
        presentations_col.find(
            {"notifications.user_id": user_id},
            {"_id": 0, "presentation_id": 1, "filename": 1, "notifications": 1},
        )
    )
    notifications = []
    for doc in docs:
        notifications.extend(build_notifications_for_user(doc, user_id))

    notifications.sort(key=lambda item: item.get("created_at") or "", reverse=True)
    unread_count = sum(1 for item in notifications if item.get("read_at") in (None, ""))
    return jsonify({"items": notifications[:20], "unread_count": unread_count})


@app.route("/api/notifications/read", methods=["POST"])
@jwt_required()
def mark_notifications_read():
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    user_id = str(current_user["_id"])
    payload = request.get_json(silent=True) or {}
    presentation_id = str(payload.get("presentation_id") or "").strip()

    query = {"notifications.user_id": user_id}
    if presentation_id:
        query["presentation_id"] = presentation_id

    presentations_col.update_many(
        query,
        {"$set": {"notifications.$[item].read_at": utc_now()}},
        array_filters=[{"item.user_id": user_id, "item.read_at": None}],
    )
    return jsonify({"ok": True})

def find_soffice() -> str | None:
    if LIBREOFFICE_PATH and Path(LIBREOFFICE_PATH).exists():
        return LIBREOFFICE_PATH
    found = shutil.which("soffice")
    if found:
        return found
    for c in [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        *glob.glob(r"C:\Program Files\LibreOffice*\program\soffice.exe"),
        "/usr/bin/soffice",
        "/usr/local/bin/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        *glob.glob("/opt/libreoffice*/program/soffice"),
    ]:
        if c and Path(c).exists():
            return c
    return None


def get_slide_count(file_path: Path) -> int:
    try:
        return len(Presentation(str(file_path)).slides)
    except Exception:
        return 0


def extract_slides_text(file_path: Path) -> list[dict]:
    """
    Extract every text shape from a PPTX.
    Returns a list usable both as Claude input and as injection blueprint.
    """
    prs = Presentation(str(file_path))
    result = []
    for slide_idx, slide in enumerate(prs.slides):
        shapes_data = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            texts = [p.text for p in shape.text_frame.paragraphs]
            if not any(t.strip() for t in texts):
                continue
            ph_idx = None
            try:
                if shape.is_placeholder:
                    ph_idx = shape.placeholder_format.idx
            except Exception:
                ph_idx = None
            shapes_data.append({
                "shape_id":        shape.shape_id,
                "shape_name":      shape.name,
                "placeholder_idx": ph_idx,
                "texts":           texts,
            })
        result.append({"slide_index": slide_idx, "shapes": shapes_data})
    return result


def extract_slide_shapes(file_path: Path, slide_index: int) -> list[dict]:
    prs = Presentation(str(file_path))
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError("Slide index out of range")

    shapes_data = []
    slide = prs.slides[slide_index]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        texts = [p.text for p in shape.text_frame.paragraphs]
        if not any(t.strip() for t in texts):
            continue
        shapes_data.append({
            "shape_id": shape.shape_id,
            "shape_name": shape.name,
            "texts": texts,
        })
    return shapes_data


def normalize_placeholder_key(key: str) -> str:
    return re.sub(r"[^a-z0-9]+", "_", str(key).strip().lower()).strip("_")


def build_placeholder_context(form_data: dict) -> dict[str, str]:
    context: dict[str, str] = {}

    for key, value in (form_data or {}).items():
        if value is None:
            continue
        if isinstance(value, list):
            string_value = ", ".join(str(item).strip() for item in value if str(item).strip())
        if isinstance(value, (str, int, float)):
            string_value = str(value).strip()
        elif not isinstance(value, list):
            string_value = str(value)
        if string_value == "":
            continue
        normalized_key = normalize_placeholder_key(key)
        if normalized_key:
            context[normalized_key] = string_value

    aliases = {
        "client": context.get("clientname", ""),
        "client_name": context.get("clientname", ""),
        "company_name": context.get("clientname", ""),
        "nom_client": context.get("clientname", ""),
        "secteur": context.get("sector", ""),
        "pays": context.get("country", ""),
        "mission": context.get("missiontype", ""),
        "mission_type": context.get("missiontype", ""),
        "normes": context.get("standards", ""),
        "contexte": context.get("context", ""),
        "objectifs": context.get("objectives", ""),
        "associe": context.get("partner", ""),
        "manager_name": context.get("manager", ""),
        "honoraires": context.get("fees", ""),
        "date_limite": context.get("deadline", ""),
        "duree": context.get("duration", ""),
    }
    for alias, value in aliases.items():
        if value:
            context[alias] = value

    client_id = form_data.get("clientId")
    if client_id:
        try:
            client_doc = clients_col.find_one({"_id": ObjectId(client_id)})
        except Exception:
            client_doc = None
        if client_doc:
            client_fields = {
                "client_company_name": client_doc.get("company_name", ""),
                "client_sector": client_doc.get("sector", ""),
                "client_country": client_doc.get("country", ""),
                "client_civility": client_doc.get("civility", ""),
                "client_responsable_name": client_doc.get("responsable_name", ""),
                "client_responsable_function": client_doc.get("responsable_function", ""),
                "client_legal_form": client_doc.get("legal_form", ""),
                "client_rccm": client_doc.get("RCCM", ""),
                "client_address": client_doc.get("address", ""),
                "client_creation_date": client_doc.get("creation_date", ""),
                "responsable_name": client_doc.get("responsable_name", ""),
                "responsable_function": client_doc.get("responsable_function", ""),
                "legal_form": client_doc.get("legal_form", ""),
                "rccm": client_doc.get("RCCM", ""),
                "address": client_doc.get("address", ""),
            }
            for key, value in client_fields.items():
                if value:
                    context[key] = str(value).strip()

    now = utc_now()
    context.setdefault("generated_date", now.strftime("%Y-%m-%d"))
    context.setdefault("generated_at", now.strftime("%Y-%m-%d %H:%M"))
    context.setdefault("current_year", str(now.year))
    return context


PLACEHOLDER_PATTERNS = [
    re.compile(r"\{\{\s*([a-zA-Z0-9_.-]+)\s*\}\}"),
    re.compile(r"\[\[\s*([a-zA-Z0-9_.-]+)\s*\]\]"),
    re.compile(r"<<\s*([a-zA-Z0-9_.-]+)\s*>>"),
    re.compile(r"%\s*([a-zA-Z0-9_.-]+)\s*%"),
]


def replace_placeholders_in_text(text: str, context: dict[str, str]) -> str:
    replaced = text
    for pattern in PLACEHOLDER_PATTERNS:
        replaced = pattern.sub(
            lambda match: context.get(normalize_placeholder_key(match.group(1)), match.group(0)),
            replaced,
        )
    return replaced


def apply_placeholder_replacements(slides_structure: list[dict], form_data: dict) -> list[dict]:
    context = build_placeholder_context(form_data)
    generated_content = []

    for slide_data in slides_structure:
        generated_shapes = []
        for shape_data in slide_data["shapes"]:
            generated_shapes.append({
                **shape_data,
                "texts": [
                    replace_placeholders_in_text(text, context)
                    for text in shape_data["texts"]
                ],
            })
        generated_content.append({
            "slide_index": slide_data["slide_index"],
            "shapes": generated_shapes,
        })

    return generated_content


def inject_content_into_pptx(template_path: Path, content: list[dict], out_file: Path):
    """
    Copy template and replace paragraph texts while preserving all formatting.
    `content` has the same shape as extract_slides_text() output.
    """
    prs = Presentation(str(template_path))

    lookup: dict[tuple, list[str]] = {}
    for slide_data in content:
        for shape_data in slide_data["shapes"]:
            lookup[(slide_data["slide_index"], shape_data["shape_id"])] = shape_data["texts"]

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            new_texts = lookup.get((slide_idx, shape.shape_id))
            if new_texts is None or not shape.has_text_frame:
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                new_text = new_texts[para_idx] if para_idx < len(new_texts) else ""
                if para.runs:
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ""

    prs.save(str(out_file))


def ensure_pdf(file_path: Path, stem: str | None = None) -> Path:
    stem = stem or file_path.stem
    pdf_path = cache_path / f"{stem}.pdf"
    if pdf_path.exists():
        return pdf_path
    soffice = find_soffice()
    if not soffice:
        raise RuntimeError("LibreOffice not found. Install it or set LIBREOFFICE_PATH.")
    tmp = cache_path / f"{stem}.pptx"
    if str(tmp) != str(file_path):
        shutil.copy2(str(file_path), str(tmp))
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(cache_path), str(tmp)],
        check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
    )
    return pdf_path


def pdf_to_images(pdf_path: Path, stem: str) -> list[str]:
    prefix = cache_path / f"{stem}_slide"
    pdftoppm = shutil.which("pdftoppm")
    if pdftoppm:
        subprocess.run(
            [pdftoppm, "-png", "-r", "120", str(pdf_path), str(prefix)],
            check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        return [base64.b64encode(p.read_bytes()).decode()
                for p in sorted(cache_path.glob(f"{stem}_slide-*.png"))]
    try:
        import fitz
        doc = fitz.open(str(pdf_path))
        return [base64.b64encode(page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5)).tobytes("png")).decode()
                for page in doc]
    except ImportError:
        pass
    raise RuntimeError("Install poppler-utils (pdftoppm) or pymupdf.")


def slides_to_images(pptx_path: Path, stem: str) -> list[str]:
    existing = sorted(cache_path.glob(f"{stem}_slide-*.png"))
    if existing:
        return [base64.b64encode(p.read_bytes()).decode() for p in existing]
    return pdf_to_images(ensure_pdf(pptx_path, stem=stem), stem)


def slides_fallback_text(file_path: Path) -> list[dict]:
    prs = Presentation(str(file_path))
    out = []
    for i, slide in enumerate(prs.slides):
        title, texts, img = "", [], None
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    ph = getattr(shape, "placeholder_format", None)
                    if not title and ph and ph.idx == 0:
                        title = t
                    else:
                        texts.append(t)
            if img is None and shape.shape_type == 13:
                try:
                    img = base64.b64encode(shape.image.blob).decode()
                except Exception:
                    pass
        out.append({"index": i + 1, "title": title, "text": "\n".join(texts), "image": img})
    return out


def invalidate_image_cache(stem: str):
    for f in cache_path.glob(f"{stem}_slide-*.png"):
        f.unlink(missing_ok=True)
    (cache_path / f"{stem}.pdf").unlink(missing_ok=True)


@app.route("/api/templates/upload", methods=["POST"])
def upload_template():
    """
    Upload one or more .pptx files to the library.
    multipart/form-data, field name: "files"
    Returns: { "uploaded": [{ filename, slide_count }], "errors": [...] }
    """
    if "files" not in request.files:
        return jsonify({"error": "No files field in request"}), 400

    lib = Path(LIBRARY_DIR)
    lib.mkdir(parents=True, exist_ok=True)

    uploaded, errors = [], []
    now = utc_now()

    for f in request.files.getlist("files"):
        if not f.filename:
            continue
        if not f.filename.lower().endswith(".pptx"):
            errors.append({"filename": f.filename, "error": "Seuls les fichiers .pptx sont acceptÃƒÆ’Ã‚Â©s"})
            continue

        safe_name = Path(f.filename).name
        dest = lib / safe_name

        if dest.exists():
            stem = safe_name[:-5]
            suffix = utc_now().strftime("%Y%m%d%H%M%S")
            safe_name = f"{stem}_{suffix}.pptx"
            dest = lib / safe_name

        try:
            f.save(str(dest))
            slide_count = get_slide_count(dest)
            templates_col.update_one(
                {"filename": safe_name},
                {
                    "$set": {
                        "filename":    safe_name,
                        "path":        str(dest),
                        "size":        dest.stat().st_size,
                        "slide_count": slide_count,
                        "updated_at":  now,
                    },
                    "$setOnInsert": {"created_at": now},
                },
                upsert=True,
            )
            uploaded.append({"filename": safe_name, "slide_count": slide_count})
        except Exception as e:
            errors.append({"filename": safe_name, "error": str(e)})

    return jsonify({"uploaded": uploaded, "errors": errors}), 201


def serialize(doc: dict) -> dict:
    if doc and "_id" in doc:
        doc["_id"] = str(doc["_id"])
    return doc


def scan_library() -> int:
    lib = Path(LIBRARY_DIR)
    if not lib.exists():
        return 0
    ops, now = [], utc_now()
    for f in lib.glob("*.pptx"):
        ops.append(UpdateOne(
            {"filename": f.name},
            {"$set": {"filename": f.name, "path": str(f), "size": f.stat().st_size,
                      "slide_count": get_slide_count(f), "updated_at": now},
             "$setOnInsert": {"created_at": now}},
            upsert=True,
        ))
    if ops:
        templates_col.bulk_write(ops)
    return len(ops)


@app.route("/api/health")
def health():
    return jsonify({
        "status": "ok",
        "libreoffice": find_soffice() or "not found",
        "pdftoppm":    shutil.which("pdftoppm") or "not found",
        "generation_mode": "placeholder-replacement",
    })


@app.route("/")
def root():
    return jsonify({"message": "Propale API v1"})


@app.route("/api/templates", methods=["GET"])
def list_templates():
    if request.args.get("scan") == "1":
        scan_library()
    return jsonify({"items": list(templates_col.find({}, {"_id": 0}).sort("filename", 1))})


@app.route("/api/templates/scan", methods=["POST"])
def scan_templates():
    return jsonify({"count": scan_library()})


@app.route("/api/templates/<path:filename>/pdf", methods=["GET"])
def template_pdf(filename: str):
    fp = Path(LIBRARY_DIR) / filename
    if not fp.exists():
        return jsonify({"error": "File not found"}), 404
    try:
        return send_file(ensure_pdf(fp), mimetype="application/pdf")
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/auth/change_password", methods=["POST"])
@jwt_required()
def change_password():
    user_id = get_jwt_identity()
    data = request.get_json() or {}

    old_password = data.get("old_password")
    new_password = data.get("new_password")

    if not old_password or not new_password:
        return jsonify({"error": "old_password and new_password are required"}), 400

    password_error = validate_password_strength(new_password)
    if password_error:
        return jsonify({"error": password_error}), 400

    try:
        from bson import ObjectId
        user = User.collection.find_one({"_id": ObjectId(user_id)}, {"password": 1})
        if not user:
            return jsonify({"error": "User not found"}), 404

        if not User.verify_password(user["password"], old_password):
            # The user is authenticated (JWT ok) but provided a wrong current password.
            # Return a validation error (not an auth/session error).
            return jsonify({"error": "Old password is incorrect"}), 400

        hashed_password = bcrypt.hashpw(new_password.encode("utf-8"), bcrypt.gensalt())
        User.collection.update_one(
            {"_id": ObjectId(user_id)},
            {"$set": {"password": hashed_password, "updated_at": utc_now()}},
        )

        return jsonify({"message": "Password updated successfully"}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/templates/<path:filename>/slides", methods=["GET"])
def template_slides(filename: str):
    fp = Path(LIBRARY_DIR) / filename
    if not fp.exists():
        return jsonify({"error": "File not found"}), 404
    mode = request.args.get("mode", "images")
    if mode == "text":
        try:
            s = slides_fallback_text(fp)
            return jsonify({"slide_count": len(s), "mode": "text", "slides": s})
        except Exception as e:
            return jsonify({"error": str(e)}), 500
    try:
        imgs = slides_to_images(fp, fp.stem)
        return jsonify({"filename": filename, "slide_count": len(imgs), "mode": "images", "slides": imgs})
    except RuntimeError as e:
        return jsonify({"error": str(e), "fallback_url": f"/api/templates/{filename}/slides?mode=text"}), 503
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/generate", methods=["POST"])
@jwt_required()
def generate_presentation():
    """
    POST /api/generate
    Body: { "template_filename": "...", "form": { clientName, sector, ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â¦ } }
    Returns: { "presentation_id": "ÃƒÂ¢Ã¢â€šÂ¬Ã‚Â¦", "slide_count": N }
    """
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    body = request.get_json(force=True)
    tmpl_file = body.get("template_filename")
    form_data = body.get("form", {})

    if not tmpl_file:
        return jsonify({"error": "template_filename is required"}), 400

    tmpl_path = Path(LIBRARY_DIR) / tmpl_file
    if not tmpl_path.exists():
        return jsonify({"error": f"Template not found: {tmpl_file}"}), 404

    try:
        structure = extract_slides_text(tmpl_path)
    except Exception as e:
        return jsonify({"error": f"Cannot read template: {e}"}), 500

    try:
        generated = apply_placeholder_replacements(structure, form_data)
    except Exception as e:
        return jsonify({"error": f"Placeholder replacement error: {e}"}), 500

    pres_id = str(uuid.uuid4())
    out_name = build_presentation_filename(form_data)
    out_file = output_path / f"{pres_id}.pptx"
    try:
        inject_content_into_pptx(tmpl_path, generated, out_file)
    except Exception as e:
        return jsonify({"error": f"PPTX creation failed: {e}"}), 500

    slide_count = get_slide_count(out_file)

    now = utc_now()
    doc = {
        "presentation_id": pres_id,
        "filename":        out_name,
        "path":            str(out_file),
        "template":        tmpl_file,
        "form":            form_data,
        "slide_count":     slide_count,
        "status":          "draft",
        "owner_user_id":   str(current_user["_id"]),
        "owner_name":      build_user_display_name(current_user),
        "current_reviewer_user_id": None,
        "current_reviewer_name": None,
        "current_reviewer_grade": None,
        "created_at":      now,
        "updated_at":      now,
    }
    result = presentations_col.insert_one(doc)

    try:
        slides_to_images(out_file, pres_id)
    except Exception:
        pass

    return jsonify({
        "presentation_id": pres_id,
        "mongo_id":        str(result.inserted_id),
        "filename":        out_name,
        "slide_count":     slide_count,
        "generation_mode": "placeholder-replacement",
    }), 201


@app.route("/api/presentations", methods=["GET"])
@jwt_required()
def list_presentations():
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    user_id = str(current_user["_id"])
    docs = list(
        presentations_col.find(
            {
                "$or": [
                    {"owner_user_id": user_id},
                    {"current_reviewer_user_id": user_id},
                    {"submission_history.submitted_to_user_id": user_id},
                ]
            }
        ).sort("created_at", -1)
    )
    return jsonify({"items": [build_presentation_response(d, current_user) for d in docs]})


@app.route("/api/presentations/<pres_id>", methods=["GET"])
@jwt_required()
def get_presentation(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403
    current_user_id = str(current_user["_id"])
    if any(
        isinstance(notification, dict)
        and str(notification.get("user_id") or "") == current_user_id
        and notification.get("read_at") is None
        for notification in (doc.get("notifications") or [])
    ):
        read_at = utc_now()
        presentations_col.update_one(
            {"presentation_id": pres_id},
            {"$set": {"notifications.$[item].read_at": read_at}},
            array_filters=[{"item.user_id": current_user_id, "item.read_at": None}],
        )
        doc = presentations_col.find_one({"presentation_id": pres_id}) or doc
    return jsonify(build_presentation_response(doc, current_user))


@app.route("/api/presentations/<pres_id>/slides", methods=["GET"])
@jwt_required()
def presentation_slides(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403
    fp = Path(doc["path"])
    if not fp.exists():
        return jsonify({"error": "PPTX file missing on disk"}), 404
    mode = request.args.get("mode", "images")
    if mode == "text":
        try:
            s = slides_fallback_text(fp)
            return jsonify({"slide_count": len(s), "mode": "text", "slides": s})
        except Exception as e:
            return jsonify({"error": str(e)}), 500
    try:
        imgs = slides_to_images(fp, pres_id)
        return jsonify({"slide_count": len(imgs), "mode": "images", "slides": imgs})
    except RuntimeError as e:
        return jsonify({"error": str(e), "fallback_url": f"/api/presentations/{pres_id}/slides?mode=text"}), 503
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/presentations/<pres_id>/slides/<int:slide_index>/shapes", methods=["GET"])
@jwt_required()
def presentation_slide_shapes(pres_id: str, slide_index: int):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403

    fp = Path(doc["path"])
    if not fp.exists():
        return jsonify({"error": "PPTX file missing on disk"}), 404

    try:
        shapes = extract_slide_shapes(fp, slide_index)
    except IndexError as e:
        return jsonify({"error": str(e)}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    return jsonify({
        "slide_index": slide_index,
        "shapes": shapes,
    })


@app.route("/api/presentations/<pres_id>/slide/<int:slide_index>", methods=["PATCH"])
@jwt_required()
def patch_slide(pres_id: str, slide_index: int):
    """
    PATCH /api/presentations/<id>/slide/<n>
    Body: { "shapes": [{ "shape_id": N, "texts": ["para1", "para2"] }] }
    Saves the PPTX and invalidates the image cache for this presentation.
    """
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403
    fp = Path(doc["path"])
    if not fp.exists():
        return jsonify({"error": "PPTX missing"}), 404

    shapes_patch = {s["shape_id"]: s["texts"] for s in request.get_json(force=True).get("shapes", [])}
    try:
        prs   = Presentation(str(fp))
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.shape_id not in shapes_patch or not shape.has_text_frame:
                continue
            new_texts = shapes_patch[shape.shape_id]
            for i, para in enumerate(shape.text_frame.paragraphs):
                t = new_texts[i] if i < len(new_texts) else ""
                if para.runs:
                    para.runs[0].text = t
                    for run in para.runs[1:]:
                        run.text = ""
        prs.save(str(fp))
    except Exception as e:
        return jsonify({"error": f"Patch failed: {e}"}), 500

    invalidate_image_cache(pres_id)
    presentations_col.update_one(
        {"presentation_id": pres_id},
        {"$set": {"updated_at": utc_now()}},
    )
    return jsonify({"ok": True})


@app.route("/api/presentations/<pres_id>", methods=["PATCH"])
@jwt_required()
def rename_presentation(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403

    payload = request.get_json(force=True) or {}
    raw_filename = str(payload.get("filename") or "").strip()
    if not raw_filename:
        return jsonify({"error": "filename is required"}), 400

    clean_name = raw_filename[:-5] if raw_filename.lower().endswith(".pptx") else raw_filename
    clean_name = re.sub(r"[\\/:*?\"<>|]+", "_", clean_name).strip(" ._")
    if not clean_name:
        return jsonify({"error": "Invalid filename"}), 400

    final_filename = f"{clean_name}.pptx"
    presentations_col.update_one(
        {"presentation_id": pres_id},
        {"$set": {"filename": final_filename, "updated_at": utc_now()}},
    )
    return jsonify({"ok": True, "filename": final_filename})


@app.route("/api/presentations/<pres_id>/comments", methods=["POST"])
@jwt_required()
def add_presentation_comment(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_comment_on_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403

    payload = request.get_json(force=True) or {}
    try:
        slide_index = int(payload.get("slide_index"))
    except (TypeError, ValueError):
        return jsonify({"error": "slide_index is required"}), 400

    comment_text = str(payload.get("text") or "").strip()
    if not comment_text:
        return jsonify({"error": "Comment text is required"}), 400

    timestamp = utc_now()
    comment = {
        "comment_id": str(uuid.uuid4()),
        "slide_index": slide_index,
        "text": comment_text,
        "author_user_id": str(current_user["_id"]),
        "author_name": build_user_display_name(current_user),
        "author_grade": current_user.get("grade"),
        "visible_to_user_ids": build_comment_visibility_user_ids(doc, current_user),
        "created_at": timestamp,
        "updated_at": timestamp,
    }
    notification_recipients = [
        user_id
        for user_id in comment["visible_to_user_ids"]
        if user_id != str(current_user["_id"])
    ]
    notifications = [
        {
            "notification_id": str(uuid.uuid4()),
            "user_id": user_id,
            "presentation_id": pres_id,
            "type": "comment_added",
            "message": f"{build_user_display_name(current_user)} a ajouté un commentaire sur la slide {slide_index + 1}.",
            "slide_index": slide_index,
            "comment_id": comment["comment_id"],
            "created_at": timestamp,
            "read_at": None,
        }
        for user_id in notification_recipients
    ]
    presentations_col.update_one(
        {"presentation_id": pres_id},
        {
            "$push": {
                "slide_comments": comment,
                "notifications": {"$each": notifications},
            },
            "$set": {"updated_at": timestamp},
        },
    )
    return jsonify({"ok": True, "comment": serialize_comment(comment)}), 201


@app.route("/api/presentations/<pres_id>/return", methods=["POST"])
@jwt_required()
def return_presentation_for_corrections(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_return_presentation_for_corrections(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403

    owner_id = str(doc.get("owner_user_id") or "")
    if not owner_id:
        return jsonify({"error": "Presentation owner not found"}), 400
    owner_user = User.collection.find_one({"_id": ObjectId(owner_id)}) if ObjectId.is_valid(owner_id) else None

    timestamp = utc_now()
    return_notification = {
        "notification_id": str(uuid.uuid4()),
        "user_id": owner_id,
        "presentation_id": pres_id,
        "type": "presentation_returned",
        "message": f"{build_user_display_name(current_user)} a renvoyé la présentation pour correction.",
        "created_at": timestamp,
        "read_at": None,
    }
    presentations_col.update_one(
        {"presentation_id": pres_id},
        {
            "$set": {
                "status": "changes_requested",
                "updated_at": timestamp,
                "current_reviewer_user_id": owner_id,
                "current_reviewer_name": doc.get("owner_name"),
                "current_reviewer_grade": owner_user.get("grade") if owner_user else None,
                "returned_by_user_id": str(current_user["_id"]),
                "returned_by_name": build_user_display_name(current_user),
            },
            "$push": {
                "notifications": return_notification,
                "submission_history": {
                    "submitted_at": timestamp,
                    "submitted_by_user_id": str(current_user["_id"]),
                    "submitted_by_name": build_user_display_name(current_user),
                    "submitted_to_user_id": owner_id,
                    "submitted_to_name": doc.get("owner_name"),
                    "submitted_to_role": "assistant",
                    "status": "changes_requested",
                    "action": "returned_for_corrections",
                }
            },
        },
    )
    return jsonify({"ok": True, "status": "changes_requested"})


@app.route("/api/presentations/<pres_id>/submit", methods=["POST"])
@jwt_required()
def submit_presentation(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403
    if not can_submit_presentation(doc, current_user):
        return jsonify({"error": "You cannot submit this presentation at this stage"}), 403

    next_target = resolve_next_submission_target(current_user, doc)
    if not next_target:
        return jsonify({"error": "This presentation cannot be submitted further"}), 400

    target_role, target_reference, allowed_grades, next_status, _label = next_target
    target_user = resolve_user_from_reference(target_reference, allowed_grades=allowed_grades)
    if not target_user:
        return jsonify({"error": f"Unable to find the selected {target_role} in active users"}), 400
    if str(target_user["_id"]) == str(current_user["_id"]):
        return jsonify({"error": "You cannot submit a presentation to yourself"}), 400

    timestamp = utc_now()
    submit_notification = {
        "notification_id": str(uuid.uuid4()),
        "user_id": str(target_user["_id"]),
        "presentation_id": pres_id,
        "type": "presentation_submitted",
        "message": f"{build_user_display_name(current_user)} vous a soumis une présentation.",
        "created_at": timestamp,
        "read_at": None,
    }
    presentations_col.update_one(
        {"presentation_id": pres_id},
        {
            "$set": {
                "status": next_status,
                "submitted_at": timestamp,
                "updated_at": timestamp,
                "current_reviewer_user_id": str(target_user["_id"]),
                "current_reviewer_name": build_user_display_name(target_user),
                "current_reviewer_grade": target_user.get("grade"),
                "submitted_by_user_id": str(current_user["_id"]),
                "submitted_by_name": build_user_display_name(current_user),
            },
            "$push": {
                "notifications": submit_notification,
                "submission_history": {
                    "submitted_at": timestamp,
                    "submitted_by_user_id": str(current_user["_id"]),
                    "submitted_by_name": build_user_display_name(current_user),
                    "submitted_to_user_id": str(target_user["_id"]),
                    "submitted_to_name": build_user_display_name(target_user),
                    "submitted_to_role": target_role,
                    "status": next_status,
                }
            },
        },
    )
    return jsonify({
        "ok": True,
        "status": next_status,
        "submitted_to": serialize_user_brief(target_user),
    })


@app.route("/api/presentations/<pres_id>/validate", methods=["POST"])
@jwt_required()
def validate_presentation(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_validate_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403

    timestamp = utc_now()
    owner_id = str(doc.get("owner_user_id") or "").strip()
    notifications = []
    recipient_ids = {owner_id, *extract_manager_user_ids(doc)}
    recipient_ids.discard(str(current_user["_id"]))
    for user_id in recipient_ids:
        if not user_id:
            continue
        notifications.append({
            "notification_id": str(uuid.uuid4()),
            "user_id": user_id,
            "presentation_id": pres_id,
            "type": "presentation_validated",
            "message": f"{build_user_display_name(current_user)} a validé la présentation.",
            "created_at": timestamp,
            "read_at": None,
        })

    presentations_col.update_one(
        {"presentation_id": pres_id},
        {
            "$set": {
                "status": "completed",
                "validated_at": timestamp,
                "validated_by_user_id": str(current_user["_id"]),
                "validated_by_name": build_user_display_name(current_user),
                "updated_at": timestamp,
                "current_reviewer_user_id": None,
                "current_reviewer_name": None,
                "current_reviewer_grade": None,
            },
            "$push": {
                "notifications": {"$each": notifications},
                "submission_history": {
                    "submitted_at": timestamp,
                    "submitted_by_user_id": str(current_user["_id"]),
                    "submitted_by_name": build_user_display_name(current_user),
                    "submitted_to_user_id": None,
                    "submitted_to_name": None,
                    "submitted_to_role": "completed",
                    "status": "completed",
                    "action": "validated",
                },
            },
        },
    )
    return jsonify({"ok": True, "status": "completed"})


@app.route("/api/presentations/<pres_id>/download", methods=["GET"])
@jwt_required()
def download_presentation(pres_id: str):
    current_user = get_current_user_doc()
    if not current_user:
        return jsonify({"error": "User not found"}), 404

    doc = presentations_col.find_one({"presentation_id": pres_id})
    if not doc:
        return jsonify({"error": "Not found"}), 404
    if not can_access_presentation(doc, current_user):
        return jsonify({"error": "Forbidden"}), 403
    fp = Path(doc["path"])
    if not fp.exists():
        return jsonify({"error": "File not found"}), 404
    name = doc.get("filename") or f"Propale_{doc['form'].get('clientName', 'client')}.pptx"
    return send_file(
        fp,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=name,
    )


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
